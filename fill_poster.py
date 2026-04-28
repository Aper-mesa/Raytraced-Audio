from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

fname = "8. Suggested-Template-for-POSTER.pptx"
prs = Presentation(fname)
slide = prs.slides[0]
shapes = {s.name: s for s in slide.shapes}

# Color palette
C_DARK   = RGBColor(0x1A, 0x1A, 0x2E)
C_ACCENT = RGBColor(0x16, 0x4B, 0x8C)
C_KEY    = RGBColor(0xC0, 0x39, 0x2B)
C_SUBHD  = RGBColor(0x2C, 0x3E, 0x50)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_GREEN  = RGBColor(0x1A, 0x7A, 0x4A)

BULLET = "\u2022  "
ARROW  = "\u25B6  "

FS_BODY   = 28
FS_SMALL  = 26
FS_BULLET = 27
FS_SUB    = 29
FS_REF    = 22

def clear_tf(shape):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    return tf

def fp(tf, runs, sb=None, sa=None):
    p = tf.paragraphs[0]
    if sb is not None: p.space_before = Pt(sb)
    if sa is not None: p.space_after  = Pt(sa)
    for (txt, bold, italic, color, fs) in runs:
        r = p.add_run(); r.text = txt
        if fs:             r.font.size  = Pt(fs)
        if bold  is not None: r.font.bold   = bold
        if italic is not None: r.font.italic = italic
        if color: r.font.color.rgb = color

def ap(tf, runs, sb=None, sa=None):
    p = tf.add_paragraph()
    if sb is not None: p.space_before = Pt(sb)
    if sa is not None: p.space_after  = Pt(sa)
    for (txt, bold, italic, color, fs) in runs:
        r = p.add_run(); r.text = txt
        if fs:             r.font.size  = Pt(fs)
        if bold  is not None: r.font.bold   = bold
        if italic is not None: r.font.italic = italic
        if color: r.font.color.rgb = color

# ── TITLE ────────────────────────────────────────────────────────────────────
tf = clear_tf(shapes["TextBox 5"])
fp(tf, [("Objective Head-Related Transfer Function Evaluation Metrics for More Consistent Perceptual Assessment", True, False, C_WHITE, 54)])

# ── AUTHORS ──────────────────────────────────────────────────────────────────
tf = clear_tf(shapes["TextBox 9"])
fp(tf, [
    ("Shanghao Zou", True, False, None, 38),
    ("    |    Supervisor: ", False, False, None, 34),
    ("Dr. Baha Ihnaini", True, False, None, 38),
])

# ── INSTITUTION ──────────────────────────────────────────────────────────────
tf = clear_tf(shapes["TextBox 10"])
fp(tf, [("Department of Computer Science and Technology, Wenzhou-Kean University, Wenzhou, China", False, True, None, 30)])

# ── INTRODUCTION ─────────────────────────────────────────────────────────────
tf = clear_tf(shapes["TextBox 12"])
fp(tf, [("What is an HRTF?", True, False, C_SUBHD, FS_SUB)], sa=4)
ap(tf, [("Head-Related Transfer Functions (HRTFs) are acoustic filters encoding how an individual's anatomy (head, pinnae, torso) transforms sound — the cornerstone of binaural spatial audio over headphones.", False, False, C_DARK, FS_BODY)], sa=10)
ap(tf, [("Why Current Metrics Fail", True, False, C_SUBHD, FS_SUB)], sa=4)
ap(tf, [(BULLET, True, False, C_ACCENT, FS_BULLET), ("LSD is magnitude-only: ", True, False, C_DARK, FS_BULLET), ("completely blind to phase distortions.", False, False, C_DARK, FS_BULLET)], sa=4)
ap(tf, [(BULLET, True, False, C_ACCENT, FS_BULLET), ("LSD is monaural: ", True, False, C_DARK, FS_BULLET), ("misses binaural cues (ILD, ITD) critical for spatial localisation.", False, False, C_DARK, FS_BULLET)], sa=4)
ap(tf, [(BULLET, True, False, C_ACCENT, FS_BULLET), ("LSD is frequency-blind: ", True, False, C_DARK, FS_BULLET), ("weights all bands equally, obscuring critical pinna notches (4-12 kHz).", False, False, C_DARK, FS_BULLET)], sa=10)
ap(tf, [("The \"LSD Blindspot\" - Key Motivation", True, False, C_KEY, FS_SUB)], sa=4)
ap(tf, [
    ("Andreopoulou & Katz (2022) showed minimum-phase reconstruction achieves ", False, False, C_DARK, FS_BODY),
    ("LSD < 1 dB", True, False, C_ACCENT, FS_BODY),
    (" (\"perfect\" by industry standards) yet caused ", False, False, C_DARK, FS_BODY),
    ("subjective ranking collapse (r ~ 0.5)", True, False, C_KEY, FS_BODY),
    (" -- inverting assessments for ~18% of HRTFs.", False, False, C_DARK, FS_BODY),
], sa=10)
ap(tf, [("Proposed Solution", True, False, C_SUBHD, FS_SUB)], sa=4)
ap(tf, [
    ("A ", False, False, C_DARK, FS_BODY),
    ("Composite Perceptual Metric (CPM)", True, False, C_ACCENT, FS_BODY),
    (" integrating LSD, ILD error, Group Delay error, and ITD error -- capturing both magnitude and phase degradations in a single weighted score.", False, False, C_DARK, FS_BODY),
])

# ── METHODS ───────────────────────────────────────────────────────────────────
tf = clear_tf(shapes["TextBox 17"])
fp(tf, [("Dataset: RIEC HRTF Database", True, False, C_SUBHD, FS_SUB)], sa=4)
ap(tf, [(BULLET, True, False, C_ACCENT, FS_BULLET), ("20 subjects ", True, False, C_DARK, FS_BULLET), ("(10M / 10F); 865 spatial positions each; 512-sample HRIRs at 48 kHz (SOFA format)", False, False, C_DARK, FS_BULLET)], sa=3)
ap(tf, [(BULLET, True, False, C_ACCENT, FS_BULLET), ("Analysis range: ", True, False, C_DARK, FS_BULLET), ("200 to 16,000 Hz", False, False, C_DARK, FS_BULLET)], sa=10)

ap(tf, [("5 Degradation Conditions", True, False, C_SUBHD, FS_SUB)], sa=4)
ap(tf, [(ARROW, True, False, C_KEY, FS_BULLET), ("A. Min-Phase Reconstruction", True, False, C_DARK, FS_BULLET), (" -- preserves magnitude, destroys phase; ITD re-inserted.", False, False, C_DARK, FS_BULLET)], sa=4)
ap(tf, [(ARROW, True, False, C_ACCENT, FS_BULLET), ("B. 1/3-Octave Smoothing", True, False, C_DARK, FS_BULLET), (" -- mild spectral loss, phase intact.", False, False, C_DARK, FS_BULLET)], sa=4)
ap(tf, [(ARROW, True, False, C_ACCENT, FS_BULLET), ("C. 1/1-Octave Smoothing", True, False, C_DARK, FS_BULLET), (" -- aggressive; obliterates pinna notches.", False, False, C_DARK, FS_BULLET)], sa=4)
ap(tf, [(ARROW, True, False, C_ACCENT, FS_BULLET), ("D. 12-bit Quantization", True, False, C_DARK, FS_BULLET), (" -- noise floor ~-72 dB; minimal distortion.", False, False, C_DARK, FS_BULLET)], sa=4)
ap(tf, [(ARROW, True, False, C_ACCENT, FS_BULLET), ("E. 8-bit Quantization", True, False, C_DARK, FS_BULLET), (" -- noise floor ~-48 dB; clearly audible.", False, False, C_DARK, FS_BULLET)], sa=10)

ap(tf, [("Four Metrics + CPM Composite", True, False, C_SUBHD, FS_SUB)], sa=4)
ap(tf, [(BULLET, True, False, C_ACCENT, FS_BULLET), ("Metric 1 -- LSD (dB): ", True, False, C_DARK, FS_BULLET), ("baseline magnitude metric.", False, False, C_DARK, FS_BULLET)], sa=3)
ap(tf, [(BULLET, True, False, C_ACCENT, FS_BULLET), ("Metric 2 -- ILD Error (dB): ", True, False, C_DARK, FS_BULLET), ("binaural level difference; JND ~1 dB.", False, False, C_DARK, FS_BULLET)], sa=3)
ap(tf, [(BULLET, True, False, C_ACCENT, FS_BULLET), ("Metric 3 -- Group Delay Error (ms): ", True, False, C_DARK, FS_BULLET), ("phase-sensitive key discriminator.", False, False, C_DARK, FS_BULLET)], sa=3)
ap(tf, [(BULLET, True, False, C_ACCENT, FS_BULLET), ("Metric 4 -- ITD Error (ms): ", True, False, C_DARK, FS_BULLET), ("binaural timing; JND ~10-30 us.", False, False, C_DARK, FS_BULLET)], sa=8)
ap(tf, [("CPM = 0.20*(LSD/2.0) + 0.25*(ILD/1.0) + 0.35*(GD/0.08) + 0.20*(ITD/0.03)", True, False, C_ACCENT, 25)], sa=3)
ap(tf, [("Weights reflect JND-based psychoacoustic thresholds. GD receives highest weight (0.35) to detect phase distortions LSD misses.", False, True, C_DARK, 24)])

# ── RESULTS ───────────────────────────────────────────────────────────────────
tf = clear_tf(shapes["TextBox 15"])
fp(tf, [("Central Finding", True, False, C_KEY, FS_SUB)], sa=4)
ap(tf, [
    ("Minimum-phase reconstruction achieves the ", False, False, C_DARK, FS_BODY),
    ("LOWEST LSD (0.030 dB)", True, False, C_GREEN, FS_BODY),
    (" of all conditions, yet produces the ", False, False, C_DARK, FS_BODY),
    ("HIGHEST CPM (14.094)", True, False, C_KEY, FS_BODY),
    (" -- over 5x the next worst. This is the LSD blindspot empirically confirmed.", False, False, C_DARK, FS_BODY),
], sa=12)

ap(tf, [("Results  (mean +/- std, n = 20 subjects)", True, False, C_SUBHD, FS_SUB)], sa=6)

# Header
ap(tf, [("Condition           LSD (dB)   ILD Err (dB)  GD Err (ms)   CPM", True, False, C_ACCENT, 24)], sa=2)
# Min-phase row highlighted
ap(tf, [
    ("Min-Phase        ", True, False, C_KEY, 24),
    ("0.030 [*]  ", True, False, C_GREEN, 24),
    ("0.058          ", False, False, C_DARK, 24),
    ("2.335 [*]     ", True, False, C_KEY, 24),
    ("14.094 [*]", True, False, C_KEY, 24),
], sa=2)
ap(tf, [("1/3-Oct Smooth   2.332       3.122          0.176          1.880", False, False, C_DARK, 24)], sa=2)
ap(tf, [("1/1-Oct Smooth   3.829       4.562          0.209          2.592", False, False, C_DARK, 24)], sa=2)
ap(tf, [("12-bit Quant.      0.140       0.238          0.060          0.352", False, False, C_DARK, 24)], sa=2)
ap(tf, [("8-bit Quant.       1.065       1.730          0.202          1.503", False, False, C_DARK, 24)], sa=4)
ap(tf, [("[*] = extreme value   |   GD weight = 0.35 (highest) -- phase-sensitive key discriminator", False, True, C_SUBHD, 22)], sa=12)

ap(tf, [("Key Observations", True, False, C_SUBHD, FS_SUB)], sa=4)
ap(tf, [
    (BULLET, True, False, C_KEY, FS_BULLET),
    ("Group Delay Error (2.335 ms) is 11-39x larger ", True, False, C_DARK, FS_BULLET),
    ("than any other condition -- directly exposing the phase destruction LSD cannot detect.", False, False, C_DARK, FS_BULLET),
], sa=5)
ap(tf, [
    (BULLET, True, False, C_ACCENT, FS_BULLET),
    ("CPM ranks all non-phase conditions monotonically ", True, False, C_DARK, FS_BULLET),
    ("(12-bit < 8-bit < 1/3-oct < 1/1-oct), confirming CPM retains full magnitude sensitivity.", False, False, C_DARK, FS_BULLET),
], sa=5)
ap(tf, [
    (BULLET, True, False, C_ACCENT, FS_BULLET),
    ("LSD vs. CPM scatter: ", True, False, C_DARK, FS_BULLET),
    ("min-phase sits at (LSD~0, CPM~14) -- far off the ideal diagonal -- while all other conditions cluster near it.", False, False, C_DARK, FS_BULLET),
], sa=5)
ap(tf, [
    (BULLET, True, False, C_ACCENT, FS_BULLET),
    ("Spectral plots (Figs. 3 & 4): ", True, False, C_DARK, FS_BULLET),
    ("magnitude spectra of original vs. min-phase are nearly identical (explaining low LSD), but group delay diverges massively in the pinna notch region (4-12 kHz).", False, False, C_DARK, FS_BULLET),
])

# ── CONCLUSIONS ───────────────────────────────────────────────────────────────
tf = clear_tf(shapes["TextBox 16"])
fp(tf, [("Three Principal Contributions", True, False, C_SUBHD, FS_SUB)], sa=6)

ap(tf, [("1.  ", True, False, C_KEY, FS_BODY), ("LSD Blindspot Empirically Confirmed", True, False, C_DARK, FS_BODY)], sa=3)
ap(tf, [("     Min-phase: LSD = 0.030 dB (best) yet CPM = 14.094 (worst). Direct empirical support for Andreopoulou & Katz (2022) on an independent database (RIEC vs. LISTEN).", False, False, C_DARK, FS_SMALL)], sa=8)

ap(tf, [("2.  ", True, False, C_ACCENT, FS_BODY), ("Group Delay Error as Phase Discriminator", True, False, C_DARK, FS_BODY)], sa=3)
ap(tf, [("     At 2.335 ms -- over 10x larger than any other condition -- GD error is the critical component that catches what LSD misses entirely.", False, False, C_DARK, FS_SMALL)], sa=8)

ap(tf, [("3.  ", True, False, C_GREEN, FS_BODY), ("CPM Captures Both Degradation Domains", True, False, C_DARK, FS_BODY)], sa=3)
ap(tf, [("     CPM correctly ranks all five conditions -- detecting phase distortions while remaining monotonically sensitive to magnitude distortions.", False, False, C_DARK, FS_SMALL)], sa=10)

ap(tf, [("Limitations", True, False, C_SUBHD, FS_SUB)], sa=4)
ap(tf, [(BULLET, True, False, C_ACCENT, FS_BULLET), ("No subjective listening tests; CPM-perception correlation unvalidated.", False, False, C_DARK, FS_SMALL)], sa=3)
ap(tf, [(BULLET, True, False, C_ACCENT, FS_BULLET), ("CPM weights are heuristic (JND-based), not optimised against subjective scores.", False, False, C_DARK, FS_SMALL)], sa=3)
ap(tf, [(BULLET, True, False, C_ACCENT, FS_BULLET), ("Only 20 of 51 available subjects used.", False, False, C_DARK, FS_SMALL)], sa=10)

ap(tf, [("Future Work", True, False, C_SUBHD, FS_SUB)], sa=4)
ap(tf, [(ARROW, True, False, C_ACCENT, FS_BULLET), ("MUSHRA subjective validation (15 expert listeners, ITU-R BS.1534-3).", False, False, C_DARK, FS_SMALL)], sa=3)
ap(tf, [(ARROW, True, False, C_ACCENT, FS_BULLET), ("SADIE II database cross-validation for generalisability.", False, False, C_DARK, FS_SMALL)], sa=3)
ap(tf, [(ARROW, True, False, C_ACCENT, FS_BULLET), ("Empirical optimisation of CPM weights via regression against MOS.", False, False, C_DARK, FS_SMALL)], sa=3)
ap(tf, [(ARROW, True, False, C_ACCENT, FS_BULLET), ("CPM as ML training loss function for HRTF generative models.", False, False, C_DARK, FS_SMALL)], sa=3)
ap(tf, [(ARROW, True, False, C_ACCENT, FS_BULLET), ("Frequency-dependent weighting to emphasise pinna notch region (4-12 kHz).", False, False, C_DARK, FS_SMALL)])

# ── REFERENCES ────────────────────────────────────────────────────────────────
tf = clear_tf(shapes["TextBox 18"])
refs = [
    "[1] Andreopoulou & Katz. Perceptual impact on localization quality evaluations of common pre-processing for non-individual HRTFs. JAES, 70(5):340-354, 2022.",
    "[2] Gebru et al. A survey of advanced spatial audio research. arXiv:2508.10924, 2025.",
    "[3] Armstrong et al. Perceptually enhanced spectral distance metric for HRTF quality prediction. JASA, 156(6):4133-4152, 2024.",
    "[4] Panah et al. BINAQUAL: A full-reference objective localization similarity metric for binaural audio. arXiv:2505.11915, 2025.",
    "[5] Ananthabhotla, Ithapu & Brimijoin. A framework for designing HRTF distance metrics that capture localization perception. JASA Express Letters, 1(4), 2021.",
    "[6] Watanabe et al. RIEC HRTF Database, Tohoku University, 2014.",
    "[7] Kistler & Wightman. A model of HRTFs based on PCA and minimum-phase reconstruction. JASA, 91(3):1637-1647, 1992.",
    "[8] Zhang et al. Towards perception-informed latent HRTF representations. arXiv:2507.02815, 2025.",
]
fp(tf, [(refs[0], False, False, C_DARK, FS_REF)], sa=4)
for ref in refs[1:]:
    ap(tf, [(ref, False, False, C_DARK, FS_REF)], sa=4)

# ── SAVE ─────────────────────────────────────────────────────────────────────
out = "POSTER-FILLED-v2.pptx"
prs.save(out)
print("Saved:", out)
